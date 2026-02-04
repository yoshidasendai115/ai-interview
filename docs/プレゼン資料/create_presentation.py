#!/usr/bin/env python3
"""
AIé¢æ¥ãƒ—ãƒ©ãƒƒãƒˆãƒ•ã‚©ãƒ¼ãƒ  æŠ•è³‡å®¶å‘ã‘ãƒ”ãƒƒãƒè³‡æ–™ç”Ÿæˆã‚¹ã‚¯ãƒªãƒ—ãƒˆ
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap


# Alias for consistency
RgbColor = RGBColor


def create_slide_with_title(prs, title_text, layout_index=1):
    """ã‚¿ã‚¤ãƒˆãƒ«ä»˜ãã‚¹ãƒ©ã‚¤ãƒ‰ã‚’ä½œæˆ"""
    slide_layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(slide_layout)
    return slide


def set_shape_fill(shape, r, g, b):
    """å›³å½¢ã®å¡—ã‚Šã¤ã¶ã—è‰²ã‚’è¨­å®š"""
    shape.fill.solid()
    shape.fill.fore_color.rgb = RgbColor(r, g, b)


def add_title_shape(slide, text, top=Inches(0.3), font_size=32):
    """ã‚¿ã‚¤ãƒˆãƒ«ãƒ†ã‚­ã‚¹ãƒˆã‚’è¿½åŠ """
    left = Inches(0.5)
    width = Inches(9)
    height = Inches(0.8)

    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0, 51, 102)  # ãƒ€ãƒ¼ã‚¯ãƒ–ãƒ«ãƒ¼
    return shape


def add_subtitle_shape(slide, text, top=Inches(1.1)):
    """ã‚µãƒ–ã‚¿ã‚¤ãƒˆãƒ«ã‚’è¿½åŠ """
    left = Inches(0.5)
    width = Inches(9)
    height = Inches(0.5)

    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(18)
    p.font.color.rgb = RgbColor(51, 102, 153)
    return shape


def add_bullet_points(slide, items, left=Inches(0.5), top=Inches(1.5), width=Inches(9), height=Inches(5)):
    """ç®‡æ¡æ›¸ãã‚’è¿½åŠ """
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # ã‚¤ãƒ³ãƒ‡ãƒ³ãƒˆãƒ¬ãƒ™ãƒ«ã®å‡¦ç†
        if isinstance(item, tuple):
            text, level = item
            p.text = "â€¢ " + text
            p.level = level
            p.font.size = Pt(16 if level == 0 else 14)
            if level > 0:
                p.font.color.rgb = RgbColor(80, 80, 80)
            else:
                p.font.color.rgb = RgbColor(51, 51, 51)
        else:
            p.text = "â€¢ " + item
            p.font.size = Pt(16)
            p.font.color.rgb = RgbColor(51, 51, 51)

        p.space_after = Pt(8)

    return shape


def add_centered_text(slide, text, top, font_size=24, bold=False, color=(51, 51, 51)):
    """ä¸­å¤®æƒãˆãƒ†ã‚­ã‚¹ãƒˆã‚’è¿½åŠ """
    left = Inches(0.5)
    width = Inches(9)
    height = Inches(1)

    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RgbColor(*color)
    return shape


def add_info_box(slide, title, items, left, top, width=Inches(4), height=Inches(2.5)):
    """æƒ…å ±ãƒœãƒƒã‚¯ã‚¹ã‚’è¿½åŠ """
    # ãƒœãƒƒã‚¯ã‚¹èƒŒæ™¯
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    set_shape_fill(box, 240, 248, 255)  # è–„ã„é’
    box.line.color.rgb = RgbColor(0, 102, 153)
    box.line.width = Pt(1)

    # ã‚¿ã‚¤ãƒˆãƒ«
    title_box = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.1), width - Inches(0.2), Inches(0.4))
    tf = title_box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0, 51, 102)

    # å†…å®¹
    content_box = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.5), width - Inches(0.2), height - Inches(0.6))
    tf = content_box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "â€¢ " + item
        p.font.size = Pt(11)
        p.font.color.rgb = RgbColor(51, 51, 51)
        p.space_after = Pt(4)

    return box


def add_header_line(slide):
    """ãƒ˜ãƒƒãƒ€ãƒ¼ãƒ©ã‚¤ãƒ³ï¼ˆé’ã„ç·šï¼‰ã‚’è¿½åŠ """
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(10), Inches(0.08)
    )
    set_shape_fill(line, 0, 102, 153)
    line.line.fill.background()


def create_presentation():
    """ãƒ—ãƒ¬ã‚¼ãƒ³ãƒ†ãƒ¼ã‚·ãƒ§ãƒ³ã‚’ä½œæˆ"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ========== ã‚¹ãƒ©ã‚¤ãƒ‰1: ã‚¿ã‚¤ãƒˆãƒ« ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # ç©ºç™½ãƒ¬ã‚¤ã‚¢ã‚¦ãƒˆ

    # èƒŒæ™¯ã®é’ã„ãƒãƒ¼
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(2.8))
    set_shape_fill(top_bar, 0, 82, 147)  # é’
    top_bar.line.fill.background()

    # mintokuãƒ­ã‚´çš„ãªãƒ†ã‚­ã‚¹ãƒˆ
    mintoku_label = add_centered_text(slide, "mintoku æ¡ç”¨æ”¯æ´ã‚µãƒ¼ãƒ“ã‚¹ æ–°æ©Ÿèƒ½", Inches(0.5), 16, False, (200, 220, 255))

    # ã‚¿ã‚¤ãƒˆãƒ«
    title = add_centered_text(slide, "AIé¢æ¥ç·´ç¿’ãƒ»è©•ä¾¡æ©Ÿèƒ½ï¼ˆä»®ï¼‰", Inches(1.1), 40, True, (255, 255, 255))

    # ã‚­ãƒ£ãƒƒãƒã‚³ãƒ”ãƒ¼
    subtitle = add_centered_text(slide, "å¤–å›½äººæ¡ç”¨ã®ã€Œè¦‹æ¥µã‚ç²¾åº¦ã€ã‚’é£›èºçš„ã«å‘ä¸Š", Inches(2), 20, False, (220, 240, 255))

    # ã‚µãƒ–æƒ…å ±
    add_centered_text(slide, "â—¯â—¯å‘ã‘æ–°ã‚µãƒ¼ãƒ“ã‚¹èª¬æ˜è³‡æ–™", Inches(4.5), 18, False, (100, 100, 100))
    add_centered_text(slide, "mintoku.com", Inches(5.5), 14, False, (0, 102, 153))
    add_centered_text(slide, "Confidential", Inches(6.5), 14, False, (150, 150, 150))

    # ========== ã‚¹ãƒ©ã‚¤ãƒ‰2: å¸‚å ´æ©Ÿä¼š ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_line(slide)
    add_title_shape(slide, "å¸‚å ´æ©Ÿä¼š")
    add_subtitle_shape(slide, "æ€¥æˆé•·ã™ã‚‹å¤–å›½äººåŠ´åƒè€…å¸‚å ´ Ã— mintokuã®é¡§å®¢åŸºç›¤")

    items = [
        "å¤–å›½äººåŠ´åƒè€…æ•°: 2024å¹´ã«200ä¸‡äººçªç ´ï¼ˆéå»æœ€é«˜ã‚’æ›´æ–°ï¼‰",
        "ç‰¹å®šæŠ€èƒ½ãƒ“ã‚¶åˆ¶åº¦ã®æ‹¡å¤§ï¼ˆ2024å¹´: å¯¾è±¡åˆ†é‡ã‚’12â†’16ã«æ‹¡å¤§ï¼‰",
        "äººæ‰‹ä¸è¶³ã®æ·±åˆ»åŒ–ï¼ˆä»‹è­·ãƒ»å»ºè¨­ãƒ»è£½é€ æ¥­ãƒ»é£²é£Ÿæ¥­ãªã©ï¼‰",
    ]
    add_bullet_points(slide, items, top=Inches(1.8), height=Inches(1.8))

    # mintokuã®å¼·ã¿
    add_info_box(slide, "mintokuã®æ—¢å­˜åŸºç›¤", [
        "å°å…¥ä¼æ¥­2,000ç¤¾ä»¥ä¸Š",
        "å¤–å›½äººæ¡ç”¨ã€œå¸°å›½ã¾ã§ä¸€å…ƒç®¡ç†",
        "Mintoku workï¼ˆæ±‚äººãƒ—ãƒ©ãƒƒãƒˆãƒ•ã‚©ãƒ¼ãƒ ï¼‰",
        "Mintoku studyï¼ˆæ—¥æœ¬èªå­¦ç¿’ã‚¢ãƒ—ãƒªï¼‰",
        "ã‚¢ã‚¸ã‚¢å„å›½ã®é€ã‚Šå‡ºã—æ©Ÿé–¢ã¨é€£æº"
    ], Inches(0.3), Inches(4), Inches(4.5), Inches(3))

    # ã‚·ãƒŠã‚¸ãƒ¼
    add_info_box(slide, "æ–°æ©Ÿèƒ½ã«ã‚ˆã‚‹ã‚·ãƒŠã‚¸ãƒ¼", [
        "æ—¢å­˜é¡§å®¢ã¸ã®è¿½åŠ ä¾¡å€¤æä¾›",
        "æ¡ç”¨ç²¾åº¦å‘ä¸Šã«ã‚ˆã‚‹é¡§å®¢æº€è¶³åº¦UP",
        "Mintoku studyã¨ã®å­¦ç¿’é€£æº",
        "æ¡ç”¨æ”¯æ´ã‚µãƒ¼ãƒ“ã‚¹ã®å·®åˆ¥åŒ–å¼·åŒ–"
    ], Inches(5.2), Inches(4), Inches(4.5), Inches(3))

    # ========== ã‚¹ãƒ©ã‚¤ãƒ‰3: èª²é¡Œæèµ· ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_line(slide)
    add_title_shape(slide, "ä¼æ¥­ãŒæŠ±ãˆã‚‹å¤–å›½äººæ¡ç”¨ã®èª²é¡Œ")

    items = [
        ("æ—¥æœ¬èªèƒ½åŠ›ã®æ­£ç¢ºãªè¦‹æ¥µã‚ãŒå›°é›£", 0),
        ("æ›¸é¡ä¸Šã®JLPTãƒ¬ãƒ™ãƒ«ï¼ˆN1ã€œN5ï¼‰ã ã‘ã§ã¯å®ŸåŠ›ãŒã‚ã‹ã‚‰ãªã„", 1),
        ("é¢æ¥ã§ã®ä¼šè©±åŠ›ã¨æ¥­å‹™é‚è¡Œèƒ½åŠ›ã«ã‚®ãƒ£ãƒƒãƒ—ãŒã‚ã‚‹", 1),
        ("ç”³å‘ŠJLPTãƒ¬ãƒ™ãƒ«ã¨å®ŸåŠ›ã®ä¹–é›¢", 0),
        ("ã€ŒN2å–å¾—ã€ã§ã‚‚å®Ÿéš›ã®ä¼šè©±ãƒ»æ•¬èªãŒä¸ååˆ†ãªã‚±ãƒ¼ã‚¹ãŒå¤šã„", 1),
        ("ãƒŸã‚¹ãƒãƒƒãƒã«ã‚ˆã‚‹æ—©æœŸé›¢è·", 0),
        ("æ¡ç”¨ã‚³ã‚¹ãƒˆæå¤±ï¼ˆ1äººã‚ãŸã‚Šå¹³å‡50ã€œ100ä¸‡å††ï¼‰", 1),
        ("ç¾å ´ã®å—ã‘å…¥ã‚Œè² æ‹…å¢—å¤§", 1),
    ]
    add_bullet_points(slide, items, top=Inches(1.5))

    # ========== ã‚¹ãƒ©ã‚¤ãƒ‰4: ã‚½ãƒªãƒ¥ãƒ¼ã‚·ãƒ§ãƒ³æ¦‚è¦ ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_line(slide)
    add_title_shape(slide, "ã‚½ãƒªãƒ¥ãƒ¼ã‚·ãƒ§ãƒ³æ¦‚è¦")
    add_subtitle_shape(slide, "AIã«ã‚ˆã‚‹é¢æ¥ç·´ç¿’ã¨ã‚ã‹ã‚Šã‚„ã™ã„è©•ä¾¡")

    # 3ã¤ã®æŸ±
    add_info_box(slide, "â‘  AIã‚¢ãƒã‚¿ãƒ¼é¢æ¥", [
        "ãƒªã‚¢ãƒ«ãªé¢æ¥å®˜ã‚¢ãƒã‚¿ãƒ¼",
        "ãƒªãƒƒãƒ—ã‚·ãƒ³ã‚¯å¯¾å¿œ",
        "JLPTãƒ¬ãƒ™ãƒ«åˆ¥ã®è³ªå•ç”Ÿæˆ",
        "ãƒªã‚¢ãƒ«ã‚¿ã‚¤ãƒ ä¼šè©±"
    ], Inches(0.3), Inches(1.8), Inches(3), Inches(2.8))

    add_info_box(slide, "â‘¡ AIè‡ªå‹•è©•ä¾¡", [
        "æ—¥æœ¬èªèƒ½åŠ›4è»¸è©•ä¾¡",
        "æ¡ç”¨é©æ€§5è»¸è©•ä¾¡",
        "JLPTãƒ¬ãƒ™ãƒ«ä¹–é›¢æ¤œå‡º",
        "æ¥­å‹™é©æ€§åˆ¤å®š"
    ], Inches(3.5), Inches(1.8), Inches(3), Inches(2.8))

    add_info_box(slide, "â‘¢ ä¼æ¥­å‘ã‘ãƒ¬ãƒãƒ¼ãƒˆ", [
        "ã‚ã‹ã‚Šã‚„ã™ã„è©•ä¾¡ãƒ‡ãƒ¼ã‚¿",
        "æ¡ç”¨åˆ¤æ–­ã®æ ¹æ‹ æä¾›",
        "æ¯”è¼ƒå¯èƒ½ãªæŒ‡æ¨™",
        "ãƒªã‚¹ã‚¯è»½æ¸›"
    ], Inches(6.7), Inches(1.8), Inches(3), Inches(2.8))

    # ä¾¡å€¤ææ¡ˆ
    value_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5), Inches(9), Inches(1.5))
    set_shape_fill(value_box, 0, 82, 147)
    value_box.line.fill.background()

    value_text = slide.shapes.add_textbox(Inches(0.5), Inches(5.3), Inches(9), Inches(1))
    tf = value_text.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "æ¡ç”¨ãƒŸã‚¹ãƒãƒƒãƒã‚’å‰Šæ¸›ã—ã€ä¼æ¥­ã®æ¡ç”¨ã‚³ã‚¹ãƒˆã¨é›¢è·ãƒªã‚¹ã‚¯ã‚’ä½æ¸›"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)

    # ========== ã‚¹ãƒ©ã‚¤ãƒ‰5: é¢æ¥ã‚·ãƒŠãƒªã‚ªæ§‹æˆ ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_line(slide)
    add_title_shape(slide, "é¢æ¥ã‚·ãƒŠãƒªã‚ªæ§‹æˆ", font_size=36)

    # ã‚¿ã‚¤ãƒ ãƒ©ã‚¤ãƒ³ãƒãƒ¼
    timeline_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2.8), Inches(9), Inches(0.15))
    set_shape_fill(timeline_bar, 200, 200, 200)
    timeline_bar.line.fill.background()

    # é¢æ¥ãƒ•ã‚§ãƒ¼ã‚ºï¼ˆã‚·ãƒ³ãƒ—ãƒ«ç‰ˆï¼‰
    phases = [
        ("å°å…¥", "5åˆ†", (100, 150, 200), 0.3),
        ("éå»", "10åˆ†", (80, 140, 100), 1.8),
        ("ç¾åœ¨", "10åˆ†", (180, 120, 60), 3.3),
        ("æœªæ¥", "10åˆ†", (140, 80, 140), 4.8),
        ("æ¡ä»¶", "5åˆ†", (120, 120, 120), 6.3),
        ("ç· ã‚", "5åˆ†", (0, 82, 147), 7.8),
    ]

    for name, time, color, left in phases:
        # å††å½¢ã®ãƒ•ã‚§ãƒ¼ã‚ºãƒãƒ¼ã‚«ãƒ¼
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left + 0.4), Inches(2.55), Inches(0.7), Inches(0.7))
        set_shape_fill(circle, *color)
        circle.line.fill.background()

        # ãƒ•ã‚§ãƒ¼ã‚ºåï¼ˆä¸Šï¼‰
        name_box = slide.shapes.add_textbox(Inches(left), Inches(1.7), Inches(1.5), Inches(0.8))
        tf = name_box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = name
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RgbColor(*color)

        # æ™‚é–“ï¼ˆå††ã®ä¸­ï¼‰
        time_box = slide.shapes.add_textbox(Inches(left + 0.4), Inches(2.65), Inches(0.7), Inches(0.5))
        tf = time_box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = time
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RgbColor(255, 255, 255)

    # çŸ¢å°ï¼ˆé€²è¡Œæ–¹å‘ï¼‰
    for x in [1.55, 3.05, 4.55, 6.05, 7.55]:
        arrow = slide.shapes.add_textbox(Inches(x), Inches(2.65), Inches(0.4), Inches(0.5))
        tf = arrow.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = "â–¶"
        p.font.size = Pt(16)
        p.font.color.rgb = RgbColor(150, 150, 150)

    # ä¸‹éƒ¨ï¼šè©•ä¾¡5è»¸ï¼ˆå¤§ããªã‚¢ã‚¤ã‚³ãƒ³é¢¨ï¼‰
    eval_title = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.6))
    tf = eval_title.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "æ¡ç”¨é©æ€§è©•ä¾¡ 5è»¸"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0, 51, 102)

    axes = [
        ("é©å¿œåŠ›", (0, 120, 180)),
        ("ã‚³ãƒŸãƒ¥ãƒ‹ã‚±ãƒ¼ã‚·ãƒ§ãƒ³", (0, 150, 100)),
        ("ä¸»ä½“æ€§", (200, 130, 50)),
        ("å®šç€æ„å‘", (150, 80, 150)),
        ("å”èª¿æ€§", (100, 100, 180)),
    ]

    for i, (name, color) in enumerate(axes):
        left = 0.5 + i * 1.9
        # å¤§ããªä¸¸
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(4.6), Inches(1.7), Inches(1.7))
        set_shape_fill(circle, *color)
        circle.line.fill.background()

        # è»¸å
        name_box = slide.shapes.add_textbox(Inches(left), Inches(5.15), Inches(1.7), Inches(0.7))
        tf = name_box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = name
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RgbColor(255, 255, 255)

    # æ³¨é‡ˆ
    note = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(9), Inches(0.5))
    tf = note.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "å„ãƒ•ã‚§ãƒ¼ã‚ºã®å›ç­”ã‹ã‚‰AIãŒ5è»¸ã‚’ç·åˆè©•ä¾¡"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(14)
    p.font.color.rgb = RgbColor(100, 100, 100)

    # ========== ã‚¹ãƒ©ã‚¤ãƒ‰6: é¢æ¥ãƒ•ãƒ­ãƒ¼ ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_line(slide)
    add_title_shape(slide, "é¢æ¥ã®æµã‚Œ", font_size=36)

    # 4ã‚¹ãƒ†ãƒƒãƒ—ã®ãƒ•ãƒ­ãƒ¼
    steps = [
        ("1", "AIãŒè³ªå•", "ã‚¢ãƒã‚¿ãƒ¼ãŒ\nè³ªå•ã‚’ç™ºè©±", (0, 82, 147)),
        ("2", "æ¬¡ã®è³ªå•ã‚’ç¢ºèª", "è³ªå•å†…å®¹\nãƒ¢ãƒ¼ãƒ€ãƒ«è¡¨ç¤º", (100, 100, 100)),
        ("3", "éŸ³å£°ã§å›ç­”", "éŒ²éŸ³ãƒ¢ãƒ¼ãƒ€ãƒ«\nãƒªã‚¢ãƒ«ã‚¿ã‚¤ãƒ æ–‡å­—èµ·ã“ã—", (0, 120, 80)),
        ("4", "æ¬¡ã®è³ªå•ã¸", "10å•\nç¹°ã‚Šè¿”ã—", (180, 100, 50)),
    ]

    for i, (num, title, desc, color) in enumerate(steps):
        left = 0.4 + i * 2.4

        # å¤§ããªå††
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left + 0.35), Inches(1.6), Inches(1.5), Inches(1.5))
        set_shape_fill(circle, *color)
        circle.line.fill.background()

        # ç•ªå·
        num_text = slide.shapes.add_textbox(Inches(left + 0.35), Inches(2.0), Inches(1.5), Inches(0.7))
        tf = num_text.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = num
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RgbColor(255, 255, 255)

        # ã‚¿ã‚¤ãƒˆãƒ«
        title_box = slide.shapes.add_textbox(Inches(left), Inches(3.2), Inches(2.2), Inches(0.5))
        tf = title_box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RgbColor(*color)

        # èª¬æ˜
        desc_box = slide.shapes.add_textbox(Inches(left), Inches(3.7), Inches(2.2), Inches(0.8))
        tf = desc_box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = desc
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.color.rgb = RgbColor(80, 80, 80)

        # çŸ¢å°
        if i < 3:
            arrow = slide.shapes.add_textbox(Inches(left + 2.0), Inches(2.1), Inches(0.5), Inches(0.5))
            tf = arrow.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = "â†’"
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.color.rgb = RgbColor(180, 180, 180)

    # ========== ã‚¹ãƒ©ã‚¤ãƒ‰7: ä¼šè©±ä¾‹ ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_line(slide)
    add_title_shape(slide, "ä¼šè©±ä¾‹", font_size=36)

    # ãƒ©ãƒ™ãƒ«
    ai_label = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.3), Inches(1.3), Inches(1.2), Inches(0.6))
    set_shape_fill(ai_label, 0, 82, 147)
    ai_label.line.fill.background()

    ai_label_text = slide.shapes.add_textbox(Inches(0.3), Inches(1.4), Inches(1.2), Inches(0.4))
    tf = ai_label_text.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "AIé¢æ¥å®˜"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)

    user_label = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(1.3), Inches(1.2), Inches(0.6))
    set_shape_fill(user_label, 0, 120, 80)
    user_label.line.fill.background()

    user_label_text = slide.shapes.add_textbox(Inches(8.5), Inches(1.4), Inches(1.2), Inches(0.4))
    tf = user_label_text.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "ãƒ¦ãƒ¼ã‚¶ãƒ¼"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)

    # ä¼šè©±ãƒ‡ãƒ¼ã‚¿
    conversations = [
        ("ai", "æ—¥æœ¬ã§åƒããŸã„ã¨æ€ã£ãŸç†ç”±ã‚’æ•™ãˆã¦ãã ã•ã„ã€‚"),
        ("user", "æ—¥æœ¬ã®æŠ€è¡“ã‚’å­¦ã³ãŸã„ã‹ã‚‰ã§ã™ã€‚æ¯å›½ã§ã‚‚æ—¥æœ¬è£½å“ã¯æœ‰åã§ã€ãšã£ã¨æ†§ã‚Œã¦ã„ã¾ã—ãŸã€‚"),
        ("ai", "5å¹´å¾Œã€ã©ã®ã‚ˆã†ã«ãªã£ã¦ã„ãŸã„ã§ã™ã‹ï¼Ÿ"),
        ("user", "æ—¥æœ¬èªã‚’ã‚‚ã£ã¨ä¸Šæ‰‹ã«ãªã£ã¦ã€ãƒªãƒ¼ãƒ€ãƒ¼ã¨ã—ã¦åƒããŸã„ã§ã™ã€‚"),
        ("ai", "æ—¥æœ¬ã§åƒãã«ã‚ãŸã£ã¦ã€ä¸å®‰ãªã“ã¨ã¯ã‚ã‚Šã¾ã™ã‹ï¼Ÿ"),
        ("user", "æ—¥æœ¬èªãŒã¾ã å®Œç’§ã§ã¯ãªã„ã®ã§å°‘ã—ä¸å®‰ã§ã™ãŒã€æ¯æ—¥å‹‰å¼·ã—ã¦ã„ã¾ã™ã€‚"),
    ]

    y_pos = 2.1
    for speaker, text in conversations:
        if speaker == "ai":
            # AIï¼ˆå·¦å´ï¼‰
            bubble = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(y_pos), Inches(5.5), Inches(0.5))
            set_shape_fill(bubble, 230, 240, 255)
            bubble.line.color.rgb = RgbColor(0, 82, 147)
            bubble.line.width = Pt(2)

            text_box = slide.shapes.add_textbox(Inches(0.5), Inches(y_pos + 0.08), Inches(5.2), Inches(0.4))
            tf = text_box.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(13)
            p.font.color.rgb = RgbColor(0, 51, 102)
        else:
            # ãƒ¦ãƒ¼ã‚¶ãƒ¼ï¼ˆå³å´ï¼‰- å¹…ã‚’åºƒã’ã¦æŠ˜ã‚Šè¿”ã—ã‚’é˜²æ­¢
            bubble = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.0), Inches(y_pos), Inches(7.7), Inches(0.5))
            set_shape_fill(bubble, 230, 255, 240)
            bubble.line.color.rgb = RgbColor(0, 120, 80)
            bubble.line.width = Pt(2)

            text_box = slide.shapes.add_textbox(Inches(2.2), Inches(y_pos + 0.08), Inches(7.4), Inches(0.4))
            tf = text_box.text_frame
            tf.word_wrap = False
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(13)
            p.font.color.rgb = RgbColor(0, 80, 50)

        y_pos += 0.6

    # ä¸‹éƒ¨ï¼šè©•ä¾¡ã‚¤ãƒ¡ãƒ¼ã‚¸
    eval_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.6), Inches(9), Inches(0.7))
    set_shape_fill(eval_box, 255, 250, 240)
    eval_box.line.color.rgb = RgbColor(200, 130, 50)
    eval_box.line.width = Pt(2)

    eval_text = slide.shapes.add_textbox(Inches(0.5), Inches(6.75), Inches(9), Inches(0.4))
    tf = eval_text.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "â†’ å„å›ç­”ã‚’AIãŒè‡ªå‹•è©•ä¾¡ï¼ˆèªå½™ãƒ»æ–‡æ³•ãƒ»å†…å®¹ãƒ»æ•¬èªï¼‰"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RgbColor(180, 100, 30)

    # ========== ã‚¹ãƒ©ã‚¤ãƒ‰8: ãƒ¬ãƒ™ãƒ«ã‚¢ãƒƒãƒ—ãƒãƒ£ãƒ¬ãƒ³ã‚¸ ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_line(slide)
    add_title_shape(slide, "ãƒ¬ãƒ™ãƒ«ã‚¢ãƒƒãƒ—ãƒãƒ£ãƒ¬ãƒ³ã‚¸", font_size=36)
    add_subtitle_shape(slide, "é«˜å¾—ç‚¹ã§ä¸Šä½ãƒ¬ãƒ™ãƒ«ã«æŒ‘æˆ¦")

    # ãƒ•ãƒ­ãƒ¼å›³ã®è¨­å®š
    flow_y = 2.2  # ãƒ•ãƒ­ãƒ¼ã®ç¸¦ä½ç½®

    # Step 1: ç”³å‘Šãƒ¬ãƒ™ãƒ« N3
    step1_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(flow_y), Inches(1.5), Inches(1.5))
    set_shape_fill(step1_circle, 0, 82, 147)  # é’ç³»
    step1_circle.line.fill.background()

    step1_text = slide.shapes.add_textbox(Inches(0.5), Inches(flow_y + 0.3), Inches(1.5), Inches(0.9))
    tf = step1_text.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "N3"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    p = tf.add_paragraph()
    p.text = "ç”³å‘Šãƒ¬ãƒ™ãƒ«"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12)
    p.font.color.rgb = RgbColor(220, 240, 255)

    # çŸ¢å°1: AIé¢æ¥å®Ÿæ–½
    arrow1 = slide.shapes.add_textbox(Inches(2.1), Inches(flow_y + 0.45), Inches(1.2), Inches(0.6))
    tf = arrow1.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "â–¶"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28)
    p.font.color.rgb = RgbColor(150, 150, 150)

    arrow1_label = slide.shapes.add_textbox(Inches(1.9), Inches(flow_y - 0.4), Inches(1.5), Inches(0.4))
    tf = arrow1_label.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "AIé¢æ¥å®Ÿæ–½"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12)
    p.font.color.rgb = RgbColor(100, 100, 100)

    # Step 2: N3é¢æ¥çµæœ
    step2_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.0), Inches(flow_y), Inches(1.8), Inches(1.5))
    set_shape_fill(step2_box, 0, 150, 100)  # ç·‘ç³»
    step2_box.line.fill.background()

    step2_text = slide.shapes.add_textbox(Inches(3.0), Inches(flow_y + 0.2), Inches(1.8), Inches(1.1))
    tf = step2_text.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "N3çµæœ"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    p = tf.add_paragraph()
    p.text = "70ç‚¹"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    p = tf.add_paragraph()
    p.text = "â˜…é«˜å¾—ç‚¹"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12)
    p.font.color.rgb = RgbColor(255, 255, 150)

    # çŸ¢å°2: ãƒãƒ£ãƒ¬ãƒ³ã‚¸è§£æ”¾
    arrow2 = slide.shapes.add_textbox(Inches(4.9), Inches(flow_y + 0.45), Inches(1.2), Inches(0.6))
    tf = arrow2.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "â–¶"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28)
    p.font.color.rgb = RgbColor(0, 150, 100)

    arrow2_label = slide.shapes.add_textbox(Inches(4.7), Inches(flow_y - 0.4), Inches(1.6), Inches(0.4))
    tf = arrow2_label.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "ãƒãƒ£ãƒ¬ãƒ³ã‚¸è§£æ”¾!"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0, 150, 100)

    # Step 3: N2ãƒãƒ£ãƒ¬ãƒ³ã‚¸
    step3_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.8), Inches(flow_y), Inches(1.5), Inches(1.5))
    set_shape_fill(step3_circle, 200, 130, 50)  # ã‚ªãƒ¬ãƒ³ã‚¸ç³»
    step3_circle.line.fill.background()

    step3_text = slide.shapes.add_textbox(Inches(5.8), Inches(flow_y + 0.3), Inches(1.5), Inches(0.9))
    tf = step3_text.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "N2"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    p = tf.add_paragraph()
    p.text = "ãƒãƒ£ãƒ¬ãƒ³ã‚¸"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12)
    p.font.color.rgb = RgbColor(255, 230, 200)

    # çŸ¢å°3: AIé¢æ¥å®Ÿæ–½
    arrow3 = slide.shapes.add_textbox(Inches(7.4), Inches(flow_y + 0.45), Inches(1.0), Inches(0.6))
    tf = arrow3.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "â–¶"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28)
    p.font.color.rgb = RgbColor(150, 150, 150)

    arrow3_label = slide.shapes.add_textbox(Inches(7.2), Inches(flow_y - 0.4), Inches(1.5), Inches(0.4))
    tf = arrow3_label.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "AIé¢æ¥å®Ÿæ–½"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12)
    p.font.color.rgb = RgbColor(100, 100, 100)

    # Step 4: N2é¢æ¥çµæœ
    step4_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(flow_y), Inches(1.5), Inches(1.5))
    set_shape_fill(step4_box, 150, 80, 150)  # ç´«ç³»
    step4_box.line.fill.background()

    step4_text = slide.shapes.add_textbox(Inches(8.0), Inches(flow_y + 0.2), Inches(1.5), Inches(1.1))
    tf = step4_text.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "N2çµæœ"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    p = tf.add_paragraph()
    p.text = "50ç‚¹"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    p = tf.add_paragraph()
    p.text = "æŒ‘æˆ¦ä¸­"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12)
    p.font.color.rgb = RgbColor(230, 200, 230)

    # ä¸‹éƒ¨èª¬æ˜ãƒœãƒƒã‚¯ã‚¹
    explain_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.3), Inches(9), Inches(2.5))
    set_shape_fill(explain_box, 255, 250, 240)  # è–„ã„ã‚ªãƒ¬ãƒ³ã‚¸
    explain_box.line.color.rgb = RgbColor(200, 150, 100)
    explain_box.line.width = Pt(2)

    # èª¬æ˜ã‚¿ã‚¤ãƒˆãƒ«
    explain_title = slide.shapes.add_textbox(Inches(0.6), Inches(4.4), Inches(8.8), Inches(0.5))
    tf = explain_title.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "ãƒ¬ãƒ™ãƒ«ã‚¢ãƒƒãƒ—ãƒãƒ£ãƒ¬ãƒ³ã‚¸ã®ä»•çµ„ã¿"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RgbColor(150, 100, 50)

    # èª¬æ˜å†…å®¹
    explain_content = slide.shapes.add_textbox(Inches(0.6), Inches(4.9), Inches(8.8), Inches(1.8))
    tf = explain_content.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    explanations = [
        ("ğŸ’¡ é«˜å¾—ç‚¹åŸºæº–:", "ç”³å‘Šãƒ¬ãƒ™ãƒ«ã§70ç‚¹ä»¥ä¸Šã‚’å–å¾—ã™ã‚‹ã¨ã€ä¸Šä½ãƒ¬ãƒ™ãƒ«ã¸ã®ãƒãƒ£ãƒ¬ãƒ³ã‚¸ãŒè§£æ”¾ã•ã‚Œã¾ã™", (0, 100, 50)),
        ("ğŸ’¡ å®ŸåŠ›æ¸¬å®š:", "å®ŸåŠ›ã‚’æ­£ç¢ºã«æ¸¬å®šã—ã€ãƒ¦ãƒ¼ã‚¶ãƒ¼ã«æœ€é©ãªãƒ¬ãƒ™ãƒ«ã‚’ç‰¹å®šã—ã¾ã™", (0, 82, 147)),
        ("ğŸ’¡ æ¡ç”¨æ™‚ã®ä¾¡å€¤:", "ä¸Šä½ãƒ¬ãƒ™ãƒ«ã§ã‚‚é«˜å¾—ç‚¹ãªã‚‰ã€ç”³å‘Šä»¥ä¸Šã®å®ŸåŠ›ãŒã‚ã‚‹ã“ã¨ã‚’è¨¼æ˜ã§ãã¾ã™", (150, 80, 150)),
    ]

    for i, (label, desc, color) in enumerate(explanations):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{label} {desc}"
        p.font.size = Pt(14)
        p.font.color.rgb = RgbColor(*color)
        p.space_after = Pt(8)

    # ========== ã‚¹ãƒ©ã‚¤ãƒ‰9: ä¸»ãªæ©Ÿèƒ½ ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_line(slide)
    add_title_shape(slide, "ä¸»ãªæ©Ÿèƒ½ï¼ˆã§ãã‚‹ã“ã¨ï¼‰")

    items = [
        ("JLPTãƒ¬ãƒ™ãƒ«åˆ¥å­¦ç¿’ãƒ—ãƒ©ãƒ³ï¼ˆN1ã€œN5å¯¾å¿œï¼‰", 0),
        ("å„ãƒ¬ãƒ™ãƒ«ã«æœ€é©åŒ–ã•ã‚ŒãŸè³ªå•ã¨è©•ä¾¡åŸºæº–", 1),
        ("AIã‚¢ãƒã‚¿ãƒ¼é¢æ¥", 0),
        ("ãƒªãƒƒãƒ—ã‚·ãƒ³ã‚¯å¯¾å¿œã®ãƒªã‚¢ãƒ«ãªé¢æ¥å®˜", 1),
        ("HeyGen APIã«ã‚ˆã‚‹é«˜å“è³ªãªã‚¢ãƒã‚¿ãƒ¼æ˜ åƒ", 1),
        ("ãã®å ´ã§ã‚¢ãƒ‰ãƒã‚¤ã‚¹", 0),
        ("éŸ³å£°èªè­˜ã«ã‚ˆã‚‹å³æ™‚æ–‡å­—èµ·ã“ã—ï¼ˆGoogle Speech-to-Textï¼‰", 1),
        ("å›ç­”ã”ã¨ã®è©•ä¾¡ã¨ã‚¢ãƒ‰ãƒã‚¤ã‚¹", 1),
        ("è‹¦æ‰‹å…‹æœãƒ«ãƒ¼ãƒ—", 0),
        ("AIãŒå¼±ç‚¹ã‚’ç‰¹å®šã—ã€é‡ç‚¹çš„ã«åå¾©ç·´ç¿’", 1),
        ("ãƒ‘ãƒ¼ã‚½ãƒŠãƒ©ã‚¤ã‚ºã•ã‚ŒãŸå­¦ç¿’ä½“é¨“", 1),
    ]
    add_bullet_points(slide, items, top=Inches(1.3), height=Inches(6))

    # ========== ã‚¹ãƒ©ã‚¤ãƒ‰6: è©•ä¾¡ã§å¾—ã‚‰ã‚Œã‚‹æƒ…å ± ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_line(slide)
    add_title_shape(slide, "è©•ä¾¡ã§å¾—ã‚‰ã‚Œã‚‹æƒ…å ±")
    add_subtitle_shape(slide, "3ã¤ã®è¦–ç‚¹ã§è©•ä¾¡")

    # æ—¥æœ¬èªèƒ½åŠ›è©•ä¾¡
    add_info_box(slide, "æ—¥æœ¬èªèƒ½åŠ›è©•ä¾¡ï¼ˆ4è»¸ï¼‰", [
        "èªå½™åŠ›: é©åˆ‡ãªå˜èªé¸æŠãƒ»èªå½™ã®è±Šå¯Œã•",
        "æ–‡æ³•: æ–‡æ³•çš„æ­£ç¢ºæ€§ãƒ»è¤‡é›‘ãªæ§‹æ–‡ã®ä½¿ç”¨",
        "å†…å®¹: è³ªå•ã¸ã®é©åˆ‡ãªå›ç­”ãƒ»è«–ç†æ€§",
        "æ•¬èª: ãƒ“ã‚¸ãƒã‚¹æ•¬èªã®æ­£ç¢ºãªä½¿ç”¨"
    ], Inches(0.3), Inches(1.8), Inches(4.5), Inches(2.3))

    # æ¡ç”¨é©æ€§è©•ä¾¡
    add_info_box(slide, "æ¡ç”¨é©æ€§è©•ä¾¡ï¼ˆ5è»¸ï¼‰", [
        "é©å¿œåŠ›: æ–°ç’°å¢ƒã¸ã®é †å¿œæ€§",
        "ã‚³ãƒŸãƒ¥ãƒ‹ã‚±ãƒ¼ã‚·ãƒ§ãƒ³åŠ›: æ„æ€ç–é€šèƒ½åŠ›",
        "ä¸»ä½“æ€§: è‡ªç™ºçš„è¡Œå‹•ãƒ»å•é¡Œè§£æ±ºæ„æ¬²",
        "å®šç€æ„å‘: é•·æœŸå‹¤ç¶šã®æ„æ€",
        "å”èª¿æ€§: ãƒãƒ¼ãƒ ãƒ¯ãƒ¼ã‚¯èƒ½åŠ›"
    ], Inches(5.2), Inches(1.8), Inches(4.5), Inches(2.3))

    # JLPTãƒ¬ãƒ™ãƒ«ä¹–é›¢æ¤œå‡º
    jlpt_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.5), Inches(9), Inches(2.3))
    set_shape_fill(jlpt_box, 255, 248, 240)  # è–„ã„ã‚ªãƒ¬ãƒ³ã‚¸
    jlpt_box.line.color.rgb = RgbColor(200, 100, 0)

    jlpt_title = slide.shapes.add_textbox(Inches(0.6), Inches(4.6), Inches(8.8), Inches(0.4))
    tf = jlpt_title.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "JLPTãƒ¬ãƒ™ãƒ«ä¹–é›¢æ¤œå‡º"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RgbColor(150, 70, 0)

    jlpt_content = slide.shapes.add_textbox(Inches(0.6), Inches(5.1), Inches(8.8), Inches(1.5))
    tf = jlpt_content.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "ç”³å‘Šãƒ¬ãƒ™ãƒ«ã¨å®ŸåŠ›ã®ä¹–é›¢ã‚’è‡ªå‹•æ¤œå‡º"
    p.font.size = Pt(14)
    p = tf.add_paragraph()
    p.text = "ä¾‹ï¼šã€Œç”³å‘ŠN3 â†’ å®ŸåŠ›N4ç›¸å½“ã€ã€Œç”³å‘ŠN2 â†’ å®ŸåŠ›N2ã€œN1ç›¸å½“ã€"
    p.font.size = Pt(13)
    p.font.color.rgb = RgbColor(100, 100, 100)
    p = tf.add_paragraph()
    p.text = "â†’ æ¡ç”¨å¾Œã®ãƒŸã‚¹ãƒãƒƒãƒãƒªã‚¹ã‚¯ã‚’äº‹å‰ã«æŠŠæ¡"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0, 100, 0)

    # ========== ã‚¹ãƒ©ã‚¤ãƒ‰7: ä¼æ¥­å‘ã‘çµ±åˆè©•ä¾¡ãƒ¬ãƒãƒ¼ãƒˆ ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_line(slide)
    add_title_shape(slide, "ä¼æ¥­å‘ã‘è©•ä¾¡ãƒ¬ãƒãƒ¼ãƒˆ")
    add_subtitle_shape(slide, "æ¡ç”¨åˆ¤æ–­ã«å¿…è¦ãªãƒ‡ãƒ¼ã‚¿ã‚’ã‚ã‹ã‚Šã‚„ã™ã")

    items = [
        ("æ¨å®šå®ŸåŠ›ãƒ¬ãƒ™ãƒ«", 0),
        ("ç”³å‘ŠN3 â†’ å®ŸåŠ›N2ç›¸å½“ ãªã©ã€å®Ÿéš›ã®èƒ½åŠ›ã‚’æ¨å®š", 1),
        ("ãƒ¬ãƒ™ãƒ«åˆ¥ã®å›ç­”å“è³ª", 0),
        ("N1ã€œN5å„ãƒ¬ãƒ™ãƒ«ã§ã®è³ªå•ã«ã©ã‚Œã ã‘ç­”ãˆã‚‰ã‚ŒãŸã‹", 1),
        ("ã©ã‚“ãªä»•äº‹ãŒã§ããã†ã‹", 0),
        ("ç°¡å˜ãªæŒ¨æ‹¶ãƒ»æ¡ˆå†…ãŒã§ãã‚‹", 1),
        ("å ±å‘Šãƒ»é€£çµ¡ãƒ»ç›¸è«‡ãŒã§ãã‚‹", 1),
        ("ãŠå®¢æ§˜å¯¾å¿œãƒ»é›»è©±å¯¾å¿œãŒã§ãã‚‹", 1),
        ("ç·åˆè©•ä¾¡ã‚¹ã‚³ã‚¢ï¼ˆ0-100ç‚¹ï¼‰", 0),
        ("æ¡ç”¨åˆ¤æ–­ã®å‚è€ƒã«ãªã‚‹ç‚¹æ•°", 1),
    ]
    add_bullet_points(slide, items, top=Inches(1.5), height=Inches(5.5))

    # ========== ã‚¹ãƒ©ã‚¤ãƒ‰8: ç«¶åˆå„ªä½æ€§ ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_line(slide)
    add_title_shape(slide, "ç«¶åˆå„ªä½æ€§ãƒ»å·®åˆ¥åŒ–ãƒã‚¤ãƒ³ãƒˆ")

    # 4ã¤ã®å¼·ã¿
    strengths = [
        ("AIã‚¢ãƒã‚¿ãƒ¼é¢æ¥", "æœ¬ç‰©ã®é¢æ¥å®˜ã®ã‚ˆã†ãª\nãƒªã‚¢ãƒ«ãªé¢æ¥ä½“é¨“", 0.3, 1.8),
        ("JLPTå…¨ãƒ¬ãƒ™ãƒ«å¯¾å¿œ", "N1ã€œN5ã™ã¹ã¦ã«å¯¾å¿œ\næ­£ç¢ºãªãƒ¬ãƒ™ãƒ«åˆ¤å®š", 5.2, 1.8),
        ("5ã¤ã®è¦–ç‚¹ã§è©•ä¾¡", "é©å¿œåŠ›ãƒ»ä¸»ä½“æ€§ãªã©\næ•°å€¤ã§ã‚ã‹ã‚Šã‚„ã™ã", 0.3, 4.3),
        ("è‹¦æ‰‹ã‚’è¦‹ã¤ã‘ã¦æ”¹å–„", "AIãŒå¼±ç‚¹ã‚’è¦‹ã¤ã‘\nç¹°ã‚Šè¿”ã—ç·´ç¿’", 5.2, 4.3),
    ]

    for title, desc, left, top in strengths:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(4.5), Inches(2))
        set_shape_fill(box, 240, 248, 255)
        box.line.color.rgb = RgbColor(0, 102, 153)

        title_shape = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.2), Inches(4.1), Inches(0.5))
        tf = title_shape.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RgbColor(0, 82, 147)

        desc_shape = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.7), Inches(4.1), Inches(1.2))
        tf = desc_shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = RgbColor(80, 80, 80)

    # ========== ã‚¹ãƒ©ã‚¤ãƒ‰9: ãƒ“ã‚¸ãƒã‚¹ãƒ¢ãƒ‡ãƒ« ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_line(slide)
    add_title_shape(slide, "ãƒ“ã‚¸ãƒã‚¹ãƒ¢ãƒ‡ãƒ«")
    add_subtitle_shape(slide, "mintokuæ¡ç”¨æ”¯æ´ã‚µãƒ¼ãƒ“ã‚¹ã®è¿½åŠ æ©Ÿèƒ½ã¨ã—ã¦æä¾›")

    # mintokuèª¬æ˜ãƒœãƒƒã‚¯ã‚¹
    mintoku_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.6), Inches(9), Inches(1.4))
    set_shape_fill(mintoku_box, 0, 82, 147)
    mintoku_box.line.fill.background()

    mintoku_text = slide.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(8.8), Inches(1.2))
    tf = mintoku_text.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "mintokuæ¡ç”¨æ”¯æ´ã‚µãƒ¼ãƒ“ã‚¹ã®é«˜ä»˜åŠ ä¾¡å€¤ã‚ªãƒ—ã‚·ãƒ§ãƒ³ã¨ã—ã¦å±•é–‹"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    p = tf.add_paragraph()
    p.text = "å°å…¥ä¼æ¥­2,000ç¤¾ä»¥ä¸Šã®æ—¢å­˜é¡§å®¢åŸºç›¤ã‚’æ´»ç”¨"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(14)
    p.font.color.rgb = RgbColor(220, 240, 255)
    p = tf.add_paragraph()
    p.text = "Mintoku studyï¼ˆæ—¥æœ¬èªå­¦ç¿’ã‚¢ãƒ—ãƒªï¼‰ã¨ã®é€£æºã§å­¦ç¿’â†’è©•ä¾¡ã®ä¸€è²«ã‚µãƒ¼ãƒ“ã‚¹"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(14)
    p.font.color.rgb = RgbColor(220, 240, 255)

    # æä¾›å½¢æ…‹
    add_info_box(slide, "æä¾›å½¢æ…‹", [
        "æ¡ç”¨æ”¯æ´ã‚µãƒ¼ãƒ“ã‚¹å¥‘ç´„ä¼æ¥­å‘ã‘ã‚ªãƒ—ã‚·ãƒ§ãƒ³",
        "Mintoku workæ±‚è·è€…å‘ã‘é¢æ¥ç·´ç¿’æ©Ÿèƒ½",
        "äººæç´¹ä»‹æ™‚ã®è©•ä¾¡ãƒ¬ãƒãƒ¼ãƒˆä»˜åŠ ä¾¡å€¤",
        "ç™»éŒ²æ”¯æ´æ©Ÿé–¢å‘ã‘ãƒ„ãƒ¼ãƒ«æä¾›"
    ], Inches(0.3), Inches(3.3), Inches(4.5), Inches(2.2))

    # åç›Šãƒ¢ãƒ‡ãƒ«
    add_info_box(slide, "åç›Šãƒ¢ãƒ‡ãƒ«", [
        "AIé¢æ¥å›æ•°ãƒ™ãƒ¼ã‚¹ã®å¾“é‡èª²é‡‘",
        "ä¼æ¥­å‘ã‘è©•ä¾¡ãƒ¬ãƒãƒ¼ãƒˆï¼ˆ1ä»¶å˜ä½ï¼‰",
        "æœˆé¡ãƒ—ãƒ¬ãƒŸã‚¢ãƒ ãƒ—ãƒ©ãƒ³",
        "ã‚¨ãƒ³ã‚¿ãƒ¼ãƒ—ãƒ©ã‚¤ã‚ºï¼ˆç„¡åˆ¶é™åˆ©ç”¨ï¼‰"
    ], Inches(5.2), Inches(3.3), Inches(4.5), Inches(2.2))

    # ã‚³ã‚¹ãƒˆæ§‹é€ 
    cost_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.7), Inches(9), Inches(1.6))
    set_shape_fill(cost_box, 250, 250, 250)
    cost_box.line.color.rgb = RgbColor(150, 150, 150)

    cost_title = slide.shapes.add_textbox(Inches(0.6), Inches(5.8), Inches(8.8), Inches(0.4))
    tf = cost_title.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "ä¸»è¦ã‚³ã‚¹ãƒˆæ§‹é€ ï¼ˆå¤–éƒ¨APIï¼‰"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RgbColor(80, 80, 80)

    cost_content = slide.shapes.add_textbox(Inches(0.6), Inches(6.2), Inches(8.8), Inches(1))
    tf = cost_content.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "HeyGen APIï¼ˆã‚¢ãƒã‚¿ãƒ¼ï¼‰ / Google STTï¼ˆéŸ³å£°èªè­˜ï¼‰ / OpenAI APIï¼ˆè©•ä¾¡ãƒ»è³ªå•ç”Ÿæˆï¼‰"
    p.font.size = Pt(12)
    p.font.color.rgb = RgbColor(80, 80, 80)

    # ========== ã‚¹ãƒ©ã‚¤ãƒ‰10: é–‹ç™ºå·¥æ•° ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_line(slide)
    add_title_shape(slide, "é–‹ç™ºå·¥æ•°")
    add_subtitle_shape(slide, "4ãƒ•ã‚§ãƒ¼ã‚ºæ§‹æˆã«ã‚ˆã‚‹æ®µéšçš„é–‹ç™º")

    # ã‚µãƒãƒªãƒ¼ãƒœãƒƒã‚¯ã‚¹
    summary_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.6), Inches(9), Inches(1.2))
    set_shape_fill(summary_box, 0, 82, 147)
    summary_box.line.fill.background()

    summary_text = slide.shapes.add_textbox(Inches(0.6), Inches(1.75), Inches(8.8), Inches(1))
    tf = summary_text.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "ç·é–‹ç™ºæœŸé–“: 14é€±é–“ï¼ˆç´„3.5ãƒ¶æœˆï¼‰  |  ç·å·¥æ•°: ç´„11.5äººæœˆï¼ˆ230äººæ—¥ï¼‰"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)
    p = tf.add_paragraph()
    p.text = "æ¨å¥¨ä½“åˆ¶: 4åï¼ˆPM/PLã€ãƒãƒƒã‚¯ã‚¨ãƒ³ãƒ‰ã€ãƒ•ãƒ­ãƒ³ãƒˆã‚¨ãƒ³ãƒ‰ã€QAï¼‰"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(14)
    p.font.color.rgb = RgbColor(220, 240, 255)

    # ãƒ•ã‚§ãƒ¼ã‚ºåˆ¥ãƒœãƒƒã‚¯ã‚¹
    phases = [
        ("Phase 1", "åŸºç›¤æ§‹ç¯‰", "4é€±é–“", "èªè¨¼åŸºç›¤ãƒ»APIãƒ»\nãƒ¦ãƒ¼ã‚¶ãƒ¼å‘ã‘ç”»é¢", 0.3, 3.1),
        ("Phase 2", "ã‚¢ãƒã‚¿ãƒ¼çµ±åˆ", "4é€±é–“", "HeyGenãƒ»éŸ³å£°èªè­˜ãƒ»\né¢æ¥ãƒ•ãƒ­ãƒ¼åˆ¶å¾¡", 5.2, 3.1),
        ("Phase 3", "è©•ä¾¡æ©Ÿèƒ½", "3é€±é–“", "GPT-4oè©•ä¾¡ãƒ»è‹¦æ‰‹åˆ†æãƒ»\nmintokué€£æºãƒ»ç®¡ç†ç”»é¢", 0.3, 5.0),
        ("Phase 4", "ãƒªãƒªãƒ¼ã‚¹æº–å‚™", "3é€±é–“", "ã‚»ã‚­ãƒ¥ãƒªãƒ†ã‚£ãƒ»è² è·ãƒ†ã‚¹ãƒˆãƒ»\næœ¬ç•ªç’°å¢ƒæ§‹ç¯‰", 5.2, 5.0),
    ]

    for phase_num, phase_name, duration, desc, left, top in phases:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(4.5), Inches(1.7))
        set_shape_fill(box, 240, 248, 255)
        box.line.color.rgb = RgbColor(0, 102, 153)

        # ãƒ•ã‚§ãƒ¼ã‚ºç•ªå·ã¨åå‰
        phase_title = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.1), Inches(4.2), Inches(0.4))
        tf = phase_title.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = f"{phase_num}: {phase_name}ï¼ˆ{duration}ï¼‰"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RgbColor(0, 82, 147)

        # èª¬æ˜
        phase_desc = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + 0.5), Inches(4.2), Inches(1.1))
        tf = phase_desc.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = RgbColor(80, 80, 80)

    # ========== ã‚¹ãƒ©ã‚¤ãƒ‰11: ã‚¬ãƒ³ãƒˆãƒãƒ£ãƒ¼ãƒˆ ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_line(slide)
    add_title_shape(slide, "é–‹ç™ºã‚¹ã‚±ã‚¸ãƒ¥ãƒ¼ãƒ«", font_size=36)

    # é€±ãƒ©ãƒ™ãƒ«ç”¨ã®è¨­å®š
    chart_left = 2.0  # ãƒãƒ£ãƒ¼ãƒˆé–‹å§‹ä½ç½®
    chart_width = 7.5  # ãƒãƒ£ãƒ¼ãƒˆå¹…
    week_width = chart_width / 14  # 1é€±é–“ã‚ãŸã‚Šã®å¹…
    bar_height = 0.6  # ãƒãƒ¼ã®é«˜ã•
    bar_gap = 0.3  # ãƒãƒ¼é–“ã®éš™é–“

    # é€±ç•ªå·ãƒ˜ãƒƒãƒ€ãƒ¼
    for week in range(1, 15):
        week_label = slide.shapes.add_textbox(
            Inches(chart_left + (week - 1) * week_width),
            Inches(1.4),
            Inches(week_width),
            Inches(0.4)
        )
        tf = week_label.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = str(week)
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(10)
        p.font.color.rgb = RgbColor(100, 100, 100)

    # é€±åŒºåˆ‡ã‚Šç·šï¼ˆç¸¦ç·šï¼‰
    for week in range(0, 15):
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(chart_left + week * week_width),
            Inches(1.8),
            Inches(0.01),
            Inches(4.5)
        )
        set_shape_fill(line, 220, 220, 220)
        line.line.fill.background()

    # ãƒ•ã‚§ãƒ¼ã‚ºåãƒ©ãƒ™ãƒ«ç”¨ã®å¹…
    label_width = 1.8

    # ãƒ•ã‚§ãƒ¼ã‚ºãƒ‡ãƒ¼ã‚¿: (åå‰, é–‹å§‹é€±, æœŸé–“é€±, è‰²)
    gantt_phases = [
        ("Phase 1\nåŸºç›¤æ§‹ç¯‰", 1, 4, (0, 120, 180)),
        ("Phase 2\nã‚¢ãƒã‚¿ãƒ¼çµ±åˆ", 5, 4, (0, 150, 100)),
        ("Phase 3\nè©•ä¾¡æ©Ÿèƒ½", 9, 3, (200, 130, 50)),
        ("Phase 4\nãƒªãƒªãƒ¼ã‚¹æº–å‚™", 12, 3, (150, 80, 150)),
    ]

    for i, (name, start_week, duration, color) in enumerate(gantt_phases):
        y_pos = 2.0 + i * (bar_height + bar_gap)

        # ãƒ•ã‚§ãƒ¼ã‚ºåãƒ©ãƒ™ãƒ«ï¼ˆå·¦å´ï¼‰
        label = slide.shapes.add_textbox(
            Inches(0.2),
            Inches(y_pos),
            Inches(label_width),
            Inches(bar_height)
        )
        tf = label.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RgbColor(*color)

        # ã‚¬ãƒ³ãƒˆãƒãƒ¼
        bar_left = chart_left + (start_week - 1) * week_width
        bar_width_val = duration * week_width
        bar = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(bar_left),
            Inches(y_pos),
            Inches(bar_width_val),
            Inches(bar_height)
        )
        set_shape_fill(bar, *color)
        bar.line.fill.background()

        # ãƒãƒ¼å†…ã®æœŸé–“ãƒ†ã‚­ã‚¹ãƒˆ
        bar_text = slide.shapes.add_textbox(
            Inches(bar_left),
            Inches(y_pos),
            Inches(bar_width_val),
            Inches(bar_height)
        )
        tf = bar_text.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = f"{duration}é€±é–“"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RgbColor(255, 255, 255)

    # æœˆãƒ©ãƒ™ãƒ«ï¼ˆä¸‹éƒ¨ï¼‰
    month_labels = [
        ("1ãƒ¶æœˆç›®", 0, 4),
        ("2ãƒ¶æœˆç›®", 4, 4),
        ("3ãƒ¶æœˆç›®", 8, 4),
        ("3.5ãƒ¶æœˆç›®", 12, 2),
    ]

    for label_text, start_week, weeks in month_labels:
        month_box = slide.shapes.add_textbox(
            Inches(chart_left + start_week * week_width),
            Inches(5.8),
            Inches(weeks * week_width),
            Inches(0.4)
        )
        tf = month_box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = label_text
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(11)
        p.font.color.rgb = RgbColor(80, 80, 80)

    # å‡¡ä¾‹ï¼ˆä¸‹éƒ¨ï¼‰
    legend_y = 6.4
    legend_items = [
        ("åŸºç›¤æ§‹ç¯‰", (0, 120, 180)),
        ("ã‚¢ãƒã‚¿ãƒ¼çµ±åˆ", (0, 150, 100)),
        ("è©•ä¾¡æ©Ÿèƒ½", (200, 130, 50)),
        ("ãƒªãƒªãƒ¼ã‚¹æº–å‚™", (150, 80, 150)),
    ]

    for i, (legend_text, color) in enumerate(legend_items):
        x_pos = 0.5 + i * 2.4

        # å‡¡ä¾‹ã®è‰²ãƒœãƒƒã‚¯ã‚¹
        legend_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x_pos),
            Inches(legend_y + 0.1),
            Inches(0.3),
            Inches(0.3)
        )
        set_shape_fill(legend_box, *color)
        legend_box.line.fill.background()

        # å‡¡ä¾‹ãƒ†ã‚­ã‚¹ãƒˆ
        legend_label = slide.shapes.add_textbox(
            Inches(x_pos + 0.4),
            Inches(legend_y),
            Inches(1.8),
            Inches(0.5)
        )
        tf = legend_label.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = legend_text
        p.font.size = Pt(11)
        p.font.color.rgb = RgbColor(60, 60, 60)

    # ç·æœŸé–“ã®æ³¨é‡ˆ
    total_note = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(0.4))
    tf = total_note.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "ç·é–‹ç™ºæœŸé–“: 14é€±é–“ï¼ˆç´„3.5ãƒ¶æœˆï¼‰"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0, 82, 147)

    # ========== ã‚¹ãƒ©ã‚¤ãƒ‰12: é–‹ç™ºè¦æ¨¡ ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_line(slide)
    add_title_shape(slide, "é–‹ç™ºè¦æ¨¡")
    add_subtitle_shape(slide, "ã‚·ã‚¹ãƒ†ãƒ æ§‹æˆã®æ¦‚è¦")

    # é–‹ç™ºè¦æ¨¡ãƒœãƒƒã‚¯ã‚¹
    add_info_box(slide, "ç”»é¢æ•°", [
        "ãƒ¦ãƒ¼ã‚¶ãƒ¼å‘ã‘: 5ç”»é¢",
        "ï¼ˆãƒ­ã‚°ã‚¤ãƒ³ã€ãƒ›ãƒ¼ãƒ ã€é¢æ¥ç·´ç¿’ã€",
        "  ãƒ•ã‚£ãƒ¼ãƒ‰ãƒãƒƒã‚¯ã€å­¦ç¿’é€²æ—ï¼‰",
        "ç®¡ç†è€…å‘ã‘: 5ç”»é¢",
        "åˆè¨ˆ: 10ç”»é¢"
    ], Inches(0.3), Inches(1.6), Inches(3), Inches(2.8))

    add_info_box(slide, "APIã‚¨ãƒ³ãƒ‰ãƒã‚¤ãƒ³ãƒˆ", [
        "èªè¨¼API: 3",
        "é¢æ¥ã‚»ãƒƒã‚·ãƒ§ãƒ³API: 5",
        "è©•ä¾¡API: 2",
        "è³ªå•API: 2",
        "ç®¡ç†è€…å‘ã‘API: 5",
        "åˆè¨ˆ: 19ã‚¨ãƒ³ãƒ‰ãƒã‚¤ãƒ³ãƒˆ"
    ], Inches(3.5), Inches(1.6), Inches(3), Inches(2.8))

    add_info_box(slide, "å¤–éƒ¨ã‚µãƒ¼ãƒ“ã‚¹é€£æº", [
        "HeyGen Streaming Avatar",
        "ï¼ˆã‚¢ãƒã‚¿ãƒ¼è¡¨ç¤ºãƒ»ç™ºè©±ï¼‰",
        "Google Cloud STT",
        "ï¼ˆéŸ³å£°èªè­˜ãƒ»æ–‡å­—èµ·ã“ã—ï¼‰",
        "mintoku work",
        "ï¼ˆSSOèªè¨¼ãƒ»çµæœé€ä¿¡ï¼‰"
    ], Inches(6.7), Inches(1.6), Inches(3), Inches(2.8))

    # å½¹å‰²åˆ¥å·¥æ•°
    role_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.7), Inches(9), Inches(2.5))
    set_shape_fill(role_box, 250, 250, 250)
    role_box.line.color.rgb = RgbColor(150, 150, 150)

    role_title = slide.shapes.add_textbox(Inches(0.6), Inches(4.8), Inches(8.8), Inches(0.4))
    tf = role_title.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "å½¹å‰²åˆ¥å·¥æ•°ã‚µãƒãƒªãƒ¼"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RgbColor(80, 80, 80)

    role_content = slide.shapes.add_textbox(Inches(0.6), Inches(5.3), Inches(8.8), Inches(1.8))
    tf = role_content.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    roles = [
        "â€¢ ãƒ•ãƒ­ãƒ³ãƒˆã‚¨ãƒ³ãƒ‰ã‚¨ãƒ³ã‚¸ãƒ‹ã‚¢: 84äººæ—¥ï¼ˆHeyGençµ±åˆãƒ»UIå®Ÿè£…ï¼‰",
        "â€¢ ãƒãƒƒã‚¯ã‚¨ãƒ³ãƒ‰ã‚¨ãƒ³ã‚¸ãƒ‹ã‚¢: 73äººæ—¥ï¼ˆAPIãƒ»è©•ä¾¡ãƒ­ã‚¸ãƒƒã‚¯ãƒ»å¤–éƒ¨é€£æºï¼‰",
        "â€¢ QA/ãƒ†ã‚¹ãƒˆã‚¨ãƒ³ã‚¸ãƒ‹ã‚¢: 65äººæ—¥ï¼ˆãƒ†ã‚¹ãƒˆè¨­è¨ˆãƒ»å®Ÿè¡Œãƒ»ã‚»ã‚­ãƒ¥ãƒªãƒ†ã‚£ï¼‰",
        "â€¢ ã‚¤ãƒ³ãƒ•ãƒ©ã‚¨ãƒ³ã‚¸ãƒ‹ã‚¢: 32äººæ—¥ï¼ˆAWSæ§‹ç¯‰ãƒ»CI/CDï¼‰",
        "â€¢ PM/PL: 24äººæ—¥ï¼ˆãƒ—ãƒ­ã‚¸ã‚§ã‚¯ãƒˆç®¡ç†ãƒ»ãƒ‰ã‚­ãƒ¥ãƒ¡ãƒ³ãƒˆï¼‰"
    ]
    for i, role in enumerate(roles):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = role
        p.font.size = Pt(12)
        p.font.color.rgb = RgbColor(60, 60, 60)

    # ========== ã‚¹ãƒ©ã‚¤ãƒ‰12: ãƒ©ãƒ³ãƒ‹ãƒ³ã‚°ã‚³ã‚¹ãƒˆ ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_line(slide)
    add_title_shape(slide, "ãƒ©ãƒ³ãƒ‹ãƒ³ã‚°ã‚³ã‚¹ãƒˆ")
    add_subtitle_shape(slide, "AIé¢æ¥å›æ•°ãƒ™ãƒ¼ã‚¹ã®å¾“é‡èª²é‡‘ãƒ¢ãƒ‡ãƒ«")

    # ã‚³ã‚¹ãƒˆã‚µãƒãƒªãƒ¼
    summary_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.6), Inches(9), Inches(1.8))
    set_shape_fill(summary_box, 240, 248, 255)
    summary_box.line.color.rgb = RgbColor(0, 102, 153)

    summary_title = slide.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(8.8), Inches(0.4))
    tf = summary_title.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "æœˆé¡ã‚³ã‚¹ãƒˆè©¦ç®—ï¼ˆ1å›ã®é¢æ¥ = 10åˆ†ã¨ã—ã¦è¨ˆç®—ï¼‰"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0, 82, 147)

    summary_content = slide.shapes.add_textbox(Inches(0.6), Inches(2.1), Inches(8.8), Inches(1.2))
    tf = summary_content.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "Proãƒ—ãƒ©ãƒ³ï¼ˆ50å›/æœˆï¼‰: ç´„43,000å††/æœˆ  |  Scaleãƒ—ãƒ©ãƒ³ï¼ˆ330å›/æœˆï¼‰: ç´„83,000å††/æœˆ"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RgbColor(51, 51, 51)
    p = tf.add_paragraph()
    p.text = "â€» AWSã‚¤ãƒ³ãƒ•ãƒ©ï¼ˆ27,000å††ï¼‰+ å¤–éƒ¨APIï¼ˆHeyGen / Google STT / OpenAI GPT-4oï¼‰"
    p.font.size = Pt(12)
    p.font.color.rgb = RgbColor(100, 100, 100)

    # ã‚³ã‚¹ãƒˆå†…è¨³ãƒœãƒƒã‚¯ã‚¹
    add_info_box(slide, "Proãƒ—ãƒ©ãƒ³ï¼ˆ50å›/æœˆï¼‰", [
        "HeyGen: 14,850å††",
        "Google Cloud STT: 900å††",
        "OpenAI GPT-4o: 113å††",
        "AWSã‚¤ãƒ³ãƒ•ãƒ©: 27,000å††",
        "åˆè¨ˆ: ç´„43,000å††/æœˆ"
    ], Inches(0.3), Inches(3.6), Inches(3), Inches(2.3))

    add_info_box(slide, "Scaleãƒ—ãƒ©ãƒ³ï¼ˆ330å›/æœˆï¼‰", [
        "HeyGen: 49,500å††",
        "Google Cloud STT: 5,940å††",
        "OpenAI GPT-4o: 743å††",
        "AWSã‚¤ãƒ³ãƒ•ãƒ©: 27,000å††",
        "åˆè¨ˆ: ç´„83,000å††/æœˆ"
    ], Inches(3.5), Inches(3.6), Inches(3), Inches(2.3))

    # Enterpriseãƒ—ãƒ©ãƒ³
    enterprise_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.7), Inches(3.6), Inches(3), Inches(2.3))
    set_shape_fill(enterprise_box, 255, 240, 240)
    enterprise_box.line.color.rgb = RgbColor(200, 80, 80)

    ent_title = slide.shapes.add_textbox(Inches(6.8), Inches(3.7), Inches(2.8), Inches(0.4))
    tf = ent_title.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "Enterpriseï¼ˆ331å›ä»¥ä¸Š/æœˆï¼‰"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RgbColor(180, 50, 50)

    ent_content = slide.shapes.add_textbox(Inches(6.8), Inches(4.15), Inches(2.8), Inches(1.6))
    tf = ent_content.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "â€¢ HeyGen: å€‹åˆ¥äº¤æ¸‰"
    p.font.size = Pt(11)
    p.font.color.rgb = RgbColor(80, 80, 80)
    p = tf.add_paragraph()
    p.text = "â€¢ STT/GPT: å¾“é‡èª²é‡‘"
    p.font.size = Pt(11)
    p.font.color.rgb = RgbColor(80, 80, 80)
    p = tf.add_paragraph()
    p.text = "â€¢ AWS: ã‚¹ã‚±ãƒ¼ãƒ«ã‚¢ãƒƒãƒ—"
    p.font.size = Pt(11)
    p.font.color.rgb = RgbColor(80, 80, 80)
    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "â€» è¦å€‹åˆ¥è¦‹ç©"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RgbColor(180, 50, 50)

    # æ³¨è¨˜
    note_text = slide.shapes.add_textbox(Inches(0.5), Inches(6.1), Inches(9), Inches(1))
    tf = note_text.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "â€» ç‚ºæ›¿: $1 = 150å††ã§è¨ˆç®—"
    p.font.size = Pt(11)
    p.font.color.rgb = RgbColor(120, 120, 120)
    p = tf.add_paragraph()
    p.text = "â€» Scaleï¼ˆ330å›/æœˆï¼‰ã‚’è¶…ãˆã‚‹åˆ©ç”¨ã«ã¯Enterpriseãƒ—ãƒ©ãƒ³ï¼ˆHeyGenå€‹åˆ¥äº¤æ¸‰ï¼‰ãŒå¿…è¦"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RgbColor(180, 50, 50)

    # ========== ã‚¹ãƒ©ã‚¤ãƒ‰13: å¹´é–“ã‚³ã‚¹ãƒˆ ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_header_line(slide)
    add_title_shape(slide, "å¹´é–“ã‚³ã‚¹ãƒˆè©¦ç®—", font_size=36)

    # Proãƒ—ãƒ©ãƒ³
    pro_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.8), Inches(4.3), Inches(3))
    set_shape_fill(pro_box, 240, 248, 255)
    pro_box.line.color.rgb = RgbColor(0, 82, 147)
    pro_box.line.width = Pt(2)

    pro_title = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.3), Inches(0.6))
    tf = pro_title.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "Proãƒ—ãƒ©ãƒ³"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0, 82, 147)

    pro_count = slide.shapes.add_textbox(Inches(0.5), Inches(2.6), Inches(4.3), Inches(0.5))
    tf = pro_count.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "600å›/å¹´"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18)
    p.font.color.rgb = RgbColor(100, 100, 100)

    pro_price = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(4.3), Inches(0.8))
    tf = pro_price.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "ç´„51ä¸‡å††/å¹´"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0, 82, 147)

    # Scaleãƒ—ãƒ©ãƒ³
    scale_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.8), Inches(4.3), Inches(3))
    set_shape_fill(scale_box, 0, 82, 147)
    scale_box.line.fill.background()

    scale_title = slide.shapes.add_textbox(Inches(5.2), Inches(2), Inches(4.3), Inches(0.6))
    tf = scale_title.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "Scaleãƒ—ãƒ©ãƒ³"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)

    scale_count = slide.shapes.add_textbox(Inches(5.2), Inches(2.6), Inches(4.3), Inches(0.5))
    tf = scale_count.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "3,960å›/å¹´"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(18)
    p.font.color.rgb = RgbColor(200, 220, 255)

    scale_price = slide.shapes.add_textbox(Inches(5.2), Inches(3.3), Inches(4.3), Inches(0.8))
    tf = scale_price.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "ç´„100ä¸‡å††/å¹´"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(255, 255, 255)

    # Enterprise
    ent_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.1), Inches(9), Inches(1.2))
    set_shape_fill(ent_box, 255, 245, 245)
    ent_box.line.color.rgb = RgbColor(180, 80, 80)
    ent_box.line.width = Pt(2)

    ent_text = slide.shapes.add_textbox(Inches(0.5), Inches(5.4), Inches(9), Inches(0.8))
    tf = ent_text.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "Enterpriseï¼ˆ3,960å›/å¹´ ä»¥ä¸Šï¼‰: å€‹åˆ¥è¦‹ç©"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RgbColor(180, 80, 80)

    # æ³¨è¨˜
    note = slide.shapes.add_textbox(Inches(0.5), Inches(6.6), Inches(9), Inches(0.5))
    tf = note.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "â€» ç‚ºæ›¿: $1 = 150å††ã§è¨ˆç®— | å¤–éƒ¨APIã¯USDå»ºã¦ï¼ˆç‚ºæ›¿å¤‰å‹•ã‚ã‚Šï¼‰"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12)
    p.font.color.rgb = RgbColor(120, 120, 120)

    # ========== ã‚¹ãƒ©ã‚¤ãƒ‰14: ãŠå•ã„åˆã‚ã› ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # èƒŒæ™¯
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5))
    set_shape_fill(bg, 0, 82, 147)
    bg.line.fill.background()

    # ã‚¿ã‚¤ãƒˆãƒ«
    add_centered_text(slide, "ãŠå•ã„åˆã‚ã› / Next Steps", Inches(1.5), 36, True, (255, 255, 255))

    # ãƒ‡ãƒ¢æ¡ˆå†…
    demo_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(3), Inches(6), Inches(2))
    set_shape_fill(demo_box, 255, 255, 255)
    demo_box.line.fill.background()

    demo_text = slide.shapes.add_textbox(Inches(2), Inches(3.3), Inches(6), Inches(0.5))
    tf = demo_text.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "ãƒ‡ãƒ¢ãƒ³ã‚¹ãƒˆãƒ¬ãƒ¼ã‚·ãƒ§ãƒ³ã®ã”æ¡ˆå†…"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0, 82, 147)

    demo_desc = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(6), Inches(0.8))
    tf = demo_desc.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "å®Ÿéš›ã®AIé¢æ¥ä½“é¨“ã¨è©•ä¾¡ãƒ¬ãƒãƒ¼ãƒˆã®\nãƒ‡ãƒ¢ã‚’ã”è¦§ã„ãŸã ã‘ã¾ã™"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(14)
    p.font.color.rgb = RgbColor(100, 100, 100)

    # Confidential
    add_centered_text(slide, "Confidential", Inches(6.5), 12, False, (180, 200, 220))

    return prs


def main():
    """ãƒ¡ã‚¤ãƒ³å‡¦ç†"""
    prs = create_presentation()
    output_path = "/Users/yoshidaseiichi/yoshidasendai/ai-interview/docs/ãƒ—ãƒ¬ã‚¼ãƒ³è³‡æ–™/AIé¢æ¥ãƒ—ãƒ©ãƒƒãƒˆãƒ•ã‚©ãƒ¼ãƒ _æŠ•è³‡å®¶å‘ã‘.pptx"
    prs.save(output_path)
    print(f"ãƒ—ãƒ¬ã‚¼ãƒ³ãƒ†ãƒ¼ã‚·ãƒ§ãƒ³ã‚’ä½œæˆã—ã¾ã—ãŸ: {output_path}")


if __name__ == "__main__":
    main()
